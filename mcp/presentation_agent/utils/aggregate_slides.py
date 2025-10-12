from pptx import Presentation
from pathlib import Path
import copy
import logging
import io

log = logging.getLogger(__name__)

def _clone_slide(dst_prs, src_slide, src_prs):
    """Clone a slide by deep-copying its XML and handling image relationships."""
    BLANK_LAYOUT_INDEX = 6  # usually "Blank"
    new_slide = dst_prs.slides.add_slide(dst_prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Copy all shapes - handle images separately
    for shp in src_slide.shapes:
        if "Picture" in shp.name:
            # For images, use the add_picture method with the image blob
            img = io.BytesIO(shp.image.blob)
            new_slide.shapes.add_picture(
                image_file=img,
                left=shp.left,
                top=shp.top,
                width=shp.width,
                height=shp.height
            )
        else:
            # For non-image shapes, deep copy the element
            el = shp.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Copy relationships (for charts, etc., but not notesSlides)
    for key, value in src_slide.part.rels.items():
        if "notesSlide" not in value.reltype:
            try:
                new_slide.part.rels.add_relationship(
                    value.reltype,
                    value._target,
                    value.rId
                )
            except:
                # If relationship already exists or can't be added, skip it
                pass

    # Copy background if present
    if src_slide.background and src_slide.background._element is not None:
        try:
            new_slide.background._element.clear()
            new_slide.background._element.extend(
                copy.deepcopy(list(src_slide.background._element))
            )
        except:
            pass

    # Copy notes (optional)
    if src_slide.has_notes_slide:
        try:
            notes = new_slide.notes_slide
            for shp in src_slide.notes_slide.shapes:
                if getattr(shp, "is_placeholder", False):
                    continue
                notes.shapes._spTree.insert_element_before(
                    copy.deepcopy(shp.element), 'p:extLst'
                )
        except:
            pass

    return new_slide

def assemble_pptx_from_directory(run_dir: Path, output_dir: Path, output_name: str) -> tuple[Path, int]:
    slide_files = sorted(
        [p for p in run_dir.iterdir() if p.is_file() and p.suffix.lower() == ".pptx"],
        key=lambda p: p.name.lower(),
    )
    if not slide_files:
        raise RuntimeError(f"No .pptx files found in {run_dir}")

    # Start with a new deck using size from the first deck
    first = Presentation(slide_files[0])
    dst = Presentation()  # empty base
    dst.slide_width = first.slide_width
    dst.slide_height = first.slide_height

    # Bring slides over
    for f in slide_files:
        log.info(f"Merging {f.name}")
        src = Presentation(f)
        # Normalize size to first deck to avoid mixed dimensions
        src.slide_width = dst.slide_width
        src.slide_height = dst.slide_height
        for s in src.slides:
            _clone_slide(dst, s, src)

    output_path = output_dir / output_name
    dst.save(output_path)
    return output_path, len(dst.slides)