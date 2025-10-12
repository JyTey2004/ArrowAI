"""
Slide drafting utilities for the presentation agent.

This implementation consumes statistical insights and artifacts produced by the
sandbox, enforces McKinsey-inspired structure (PEEL), and generates a PPTX
artifact via the shared sandbox.
"""

from __future__ import annotations

import json
import logging
import re
from typing import Dict, Iterable, List, Optional

from components.executor import CodeSandbox, ExecRequest

from utils import LLMAdapter
from components.models import (
    SlideArtifact,
    SlideDraft,
    SlideDraftRequest,
    SlideDraftResponse,
)

log = logging.getLogger(__name__)


SLIDE_PLANNER_PROMPT = """You are a senior engagement manager preparing a McKinsey-style slide.
Use the PEEL structure: Point (headline), Evidence (chart/callout), Explanation (2-4 bullets), Link (source footnote).

Guidelines:
- Titles are conclusions written in active voice (10-18 words).
- One slide, one job: deliver a single clear claim.
- Highlight quantitative proof using the supplied artifacts; prefer one primary chart.
- Keep total words on slide under 120; footnotes may extend slightly.
- If content is too rich for one slide, include a `followup_hint` explaining what should spill to the next slide.
- Provide JSON only with this schema:
  {{
    "title": str,
    "point": str,
    "callout": str,
    "explanation_bullets": [str, ...],
    "footnote": str,
    "recommended_artifact": Optional[str],
    "needs_followup": bool,
    "followup_hint": Optional[str]
  }}
Do NOT output markdown or commentary—JSON only."""


class SlideDraftWriter:
    """Expand outline sections into slide-level drafts using the sandbox."""

    def __init__(
        self,
        sandbox: CodeSandbox,
        llm: LLMAdapter,
        *,
        max_words: int = 120,
    ) -> None:
        self.sandbox = sandbox
        self.llm = llm
        self.max_words = max_words

    def draft_slide(self, request: SlideDraftRequest) -> SlideDraftResponse:
        """Generate a PEEL slide and corresponding PPTX artifact."""
        plan = self._plan_slide_content(request)
        bullets = plan.get("explanation_bullets") or ["Insert supporting bullet."]
        bullets = [self._normalize_bullet(b) for b in bullets][:5]

        word_total = self._word_count([plan.get("point", ""), *bullets, plan.get("footnote", "")])
        needs_followup_flag = bool(plan.get("needs_followup"))
        if word_total > self.max_words or len(bullets) > 4:
            needs_followup_flag = True
            plan.setdefault(
                "followup_hint",
                "Continue the evidence on the next slide to keep this one lean.",
            )

        ppt_artifact_path = self._render_slide_via_sandbox(request, plan, bullets, needs_followup_flag)

        slide = SlideDraft(
            title=plan.get("title") or request.section_title,
            bullets=bullets,
            speaker_notes=plan.get("footnote") if request.include_speaker_notes else None,
            suggested_visual=plan.get("recommended_artifact"),
        )

        rationale = (
            f"Slide drafted from insight summary and {len(request.artifacts)} artifact(s). "
            "PPTX artifact saved for downstream assembly."
        )
        if needs_followup_flag:
            rationale += " Additional content flagged for next slide."

        return SlideDraftResponse(
            slide=slide,
            rationale=rationale,
            needs_followup=needs_followup_flag,
            followup_hint=plan.get("followup_hint"),
            ppt_artifact_path=ppt_artifact_path,
        )

    # ------------------------------------------------------------------ #
    # Planning helpers
    # ------------------------------------------------------------------ #
    def _plan_slide_content(self, request: SlideDraftRequest) -> Dict[str, any]:
        prompt = self._build_planner_prompt(request)
        raw = self.llm.generate(prompt)
        plan = self._parse_plan_json(raw)
        if not plan.get("title"):
            plan["title"] = self._fallback_title(request)
        if not plan.get("callout"):
            plan["callout"] = self._fallback_callout(request)
        return plan

    def _build_planner_prompt(self, request: SlideDraftRequest) -> str:
        key_points_text = "\n".join(f"- {kp}" for kp in request.key_points) or "- (none provided)"
        artifact_lines = self._summarize_artifacts(request.artifacts)
        return (
            f"{SLIDE_PLANNER_PROMPT}\n\n"
            f"Section: {request.section_title}\n"
            f"Tone: {request.tone}\n"
            f"Insight summary:\n{request.insight_summary}\n\n"
            f"Key points to cover:\n{key_points_text}\n\n"
            f"Available artifacts:\n{artifact_lines}\n"
        )

    @staticmethod
    def _summarize_artifacts(artifacts: Iterable[SlideArtifact]) -> str:
        lines: List[str] = []
        for idx, art in enumerate(artifacts or [], start=1):
            descriptor = art.description or "No description"
            art_type = art.artifact_type or "unknown"
            origin = art.path or "n/a"
            primary = "PRIMARY" if art.is_primary else "supporting"
            lines.append(f"{idx}. {art.name} [{art_type}, {primary}] — {descriptor}. Source: {origin}")
        return "\n".join(lines) if lines else "None"

    @staticmethod
    def _parse_plan_json(raw: str) -> Dict[str, any]:
        text = raw.strip()
        start, end = text.find("{"), text.rfind("}")
        if start == -1 or end == -1 or end <= start:
            raise ValueError("Slide planner did not return JSON.")
        payload = json.loads(text[start : end + 1])
        payload.setdefault("explanation_bullets", [])
        payload.setdefault("footnote", "")
        return payload

    @staticmethod
    def _fallback_title(request: SlideDraftRequest) -> str:
        if request.key_points:
            return request.key_points[0][:80]
        return f"{request.section_title}: key takeaway"

    @staticmethod
    def _fallback_callout(request: SlideDraftRequest) -> str:
        return "Highlight principal metric from analysis."

    @staticmethod
    def _normalize_bullet(text: str) -> str:
        cleaned = re.sub(r"\s+", " ", text or "").strip()
        return cleaned or "Fill in supporting detail."

    @staticmethod
    def _word_count(chunks: Iterable[str]) -> int:
        return sum(len((chunk or "").split()) for chunk in chunks)

    # ------------------------------------------------------------------ #
    # Sandbox rendering
    # ------------------------------------------------------------------ #
    def _render_slide_via_sandbox(
        self,
        request: SlideDraftRequest,
        plan: Dict[str, any],
        bullets: List[str],
        needs_followup: bool,
    ) -> Optional[str]:
        slide_id = request.slide_identifier or self._slugify(plan.get("title") or request.section_title)
        filename = f"{slide_id}.pptx"
        content = {
            "title": plan.get("title") or request.section_title,
            "point": plan.get("point") or plan.get("title") or "",
            "bullets": bullets,
            "callout": plan.get("callout") or "",
            "footnote": plan.get("footnote") or "",
            "followup_hint": plan.get("followup_hint") if needs_followup else "",
            "needs_followup": needs_followup,
            "artifacts": [art.dict() for art in request.artifacts],
            "recommended_artifact": plan.get("recommended_artifact"),
            "filename": filename,
        }

        code = self._build_pptx_code(content)
        exec_request = ExecRequest(
            code=code,
            language="python",
            timeout_s=180,
            use_llm_writer=False,
        )
        result = self.sandbox.exec_cell(
            thread_id=request.thread_id,
            req=exec_request,
            code_llm=None,
            eval_llm=None,
            execution_context=request.insight_summary,
        )
        if not result.ok:
            raise RuntimeError(f"Slide rendering failed: {result.stderr}")

        ppt_path = None
        for file_obj in result.files_out or []:
            path = getattr(file_obj, "path", None)
            if path and path.lower().endswith(".pptx"):
                ppt_path = path
                break

        if not ppt_path:
            log.warning("Sandbox did not report a PPTX artifact for slide '%s'.", filename)
        return ppt_path

    @staticmethod
    def _build_pptx_code(content: Dict[str, any]) -> str:
        payload = json.dumps(content)
        return f'''
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

content = json.loads({repr(payload)})

OUTPUTS_DIR = Path(OUTPUTS_DIR)
OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
deck_path = OUTPUTS_DIR / content["filename"]

prs = Presentation()
slide_layout = prs.slide_layouts[5]  # title only
slide = prs.slides.add_slide(slide_layout)

title_shape = slide.shapes.title
title_shape.text = content["title"]
title_shape.text_frame.paragraphs[0].font.size = Pt(28)

body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.5), Inches(3.2))
body_tf = body_box.text_frame
body_tf.word_wrap = True
body_tf.clear()
body_tf.text = content["point"]
body_tf.paragraphs[0].font.size = Pt(20)

for bullet in content["bullets"]:
    p = body_tf.add_paragraph()
    p.text = bullet
    p.level = 1
    p.font.size = Pt(16)

callout_text = content.get("callout") or ""
visual_note_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.1), Inches(2.2), Inches(1.4))
visual_tf = visual_note_box.text_frame
visual_tf.word_wrap = True
visual_tf.text = callout_text
visual_tf.paragraphs[0].font.size = Pt(16)
visual_note_box.fill.solid()
visual_note_box.fill.fore_color.rgb = prs.slide_master.background.fill.fore_color.rgb

# Attempt to embed recommended artifact if local image path is available.
import os
import mimetypes

def _find_primary_artifact():
    rec = content.get("recommended_artifact")
    pool = content.get("artifacts") or []
    if rec:
        for art in pool:
            if art.get("name") == rec or art.get("path") == rec:
                return art
    for art in pool:
        if art.get("is_primary"):
            return art
    return pool[0] if pool else None

primary = _find_primary_artifact()
image_added = False
if primary:
    path = primary.get("path")
    if path and not path.startswith("s3://"):
        candidate = Path(path)
        if candidate.exists():
            mime, _ = mimetypes.guess_type(candidate.name)
            if mime and mime.startswith("image"):
                slide.shapes.add_picture(str(candidate), Inches(0.6), Inches(3.1), height=Inches(2.8))
                image_added = True

if not image_added and primary:
    banner = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(5.6), Inches(0.8))
    btf = banner.text_frame
    btf.word_wrap = True
    btf.text = f"Use artifact: {{primary.get('name')}} (source: {{primary.get('path') or 'n/a'}})"
    btf.paragraphs[0].font.size = Pt(14)

footnote_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.7))
foot_tf = footnote_box.text_frame
foot_tf.word_wrap = True
foot_tf.text = content.get("footnote") or ""
foot_tf.paragraphs[0].font.size = Pt(12)

if content.get("followup_hint"):
    hint_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(8.8), Inches(0.7))
    hint_tf = hint_box.text_frame
    hint_tf.text = f"Next slide: {{content['followup_hint']}}"
    hint_tf.paragraphs[0].font.size = Pt(12)
    hint_tf.paragraphs[0].font.italic = True

prs.save(deck_path)
print(f"EVIDENCE: key=slide_title value={{content['title']}}")
print(f"EVIDENCE: key=primary_visual value={{content.get('recommended_artifact') or 'n/a'}}")
print(f"ARTIFACT: outputs/{{content['filename']}}")
print("DONE")
'''

    @staticmethod
    def _slugify(text: str) -> str:
        slug = re.sub(r"[^a-zA-Z0-9]+", "_", (text or "slide")).strip("_").lower()
        return slug or "slide"
